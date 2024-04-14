from pptx import Presentation
from pptx.util import Inches
from pptx.enum.shapes import MSO_SHAPE

file_path = "samples/output/KittenKaboodlePresentation.pptx"
# Create a presentation object
prs = Presentation()

# Add a slide; the layout 1 is typically a title slide
slide = prs.slides.add_slide(prs.slide_layouts[5])  # Using a blank slide layout

# Define width, height, and margins
width = Inches(2.0)
height = Inches(1.0)
left = Inches(2.0)
top = Inches(2.0)

# Add a rectangle shape with text
shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
shape.text = "Kitten Kaboodle"

# Set text frame alignment (optional, to center the text in the shape)
text_frame = shape.text_frame
text_frame.text = "Kitten Kaboodle"
text_frame.paragraphs[0].alignment = 1  # 1 is for center alignment

# Save the presentation
prs.save(file_path)
